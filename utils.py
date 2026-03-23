
# utils.py   
 
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation





## bfix   
from pydub import AudioSegment
import librosa
import numpy as np
import re
## bfix   



def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pdf(path):
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_ppt(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text




## bfix   

#  Minimum duration check


def check_audio_duration(file_path: str, min_seconds: int = 2):

    audio = AudioSegment.from_file(file_path)

    duration = len(audio) / 1000

    return duration >= min_seconds, duration



#  Silence detection


def detect_silence(file_path: str):

    y, sr = librosa.load(file_path)

    energy = np.mean(np.abs(y))

    return energy






#  Transcript validation


def validate_transcript(text: str):

    if not text:
        return False

    text = text.strip()

    if text == "EMPTY_AUDIO":
        return False

    # Minimum characters
    if len(text) < 25:
        return False

    # Minimum words
    words = re.findall(r'\w+', text)

    if len(words) < 6:
        return False

    # Reject generic hallucination phrases
    banned_phrases = [
        "thank you",
        "thanks for watching",
        "background noise",
        "music playing",
        "no clear audioble voice"
    ]

    lower_text = text.lower()

    for phrase in banned_phrases:
        if phrase in lower_text:
            return False

    return True


